#!/usr/bin/env python3
"""Build a PowerPoint deck from slide images.

Each image is inserted as a full-bleed slide. Images are center-cropped to the
chosen aspect ratio first, so the final PPT never stretches/squashes artwork.
"""
from __future__ import annotations

import argparse
import re
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

SUPPORTED_IMAGES = {".png", ".jpg", ".jpeg", ".webp"}
ASPECT_PRESETS = {
    "16:9": (16, 9, 13.333333, 7.5),
    "4:3": (4, 3, 10.0, 7.5),
}


def is_metadata_sidecar(path: Path) -> bool:
    return path.name.startswith("._") or path.name == ".DS_Store" or "__MACOSX" in path.parts


def natural_key(path: Path) -> list[object]:
    return [int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", path.name)]


def list_images(images_dir: Path) -> list[Path]:
    files = [
        p
        for p in images_dir.iterdir()
        if p.is_file() and p.suffix.lower() in SUPPORTED_IMAGES and not is_metadata_sidecar(p)
    ]
    return sorted(files, key=natural_key)


def center_crop_to_aspect(src: Path, dst: Path, aspect_w: int, aspect_h: int) -> None:
    """Create a PNG cropped to target aspect without changing visual proportions."""
    with Image.open(src) as im:
        im = im.convert("RGB")
        w, h = im.size
        target = aspect_w / aspect_h
        current = w / h
        if abs(current - target) < 0.002:
            im.save(dst, "PNG")
            return
        if current > target:
            new_w = int(h * target)
            left = (w - new_w) // 2
            box = (left, 0, left + new_w, h)
        else:
            new_h = int(w / target)
            top = (h - new_h) // 2
            box = (0, top, w, top + new_h)
        im.crop(box).save(dst, "PNG")


def main() -> int:
    parser = argparse.ArgumentParser(description="Create a full-image PPTX from generated slide images.")
    parser.add_argument("images_dir", type=Path, help="Folder containing slide_01.png, slide_02.png, ...")
    parser.add_argument("--out", type=Path, default=Path("output/deck.pptx"))
    parser.add_argument("--aspect", choices=sorted(ASPECT_PRESETS), default="16:9")
    parser.add_argument("--title", default="Image PPT")
    args = parser.parse_args()

    if not args.images_dir.exists():
        raise SystemExit(f"Images folder not found: {args.images_dir}")

    images = list_images(args.images_dir)
    if not images:
        raise SystemExit(f"No images found in {args.images_dir}")

    aspect_w, aspect_h, slide_w_in, slide_h_in = ASPECT_PRESETS[args.aspect]
    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)
    blank_layout = prs.slide_layouts[6]
    prs.core_properties.title = args.title

    args.out.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="image_pptx_") as tmp:
        tmp_dir = Path(tmp)
        for idx, image_path in enumerate(images, start=1):
            prepared = tmp_dir / f"slide_{idx:03d}.png"
            center_crop_to_aspect(image_path, prepared, aspect_w, aspect_h)
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(prepared), 0, 0, width=prs.slide_width, height=prs.slide_height)
        prs.save(args.out)

    print(f"Saved PPTX: {args.out} ({len(images)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
